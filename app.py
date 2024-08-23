from flask import Flask, render_template, request, send_file, make_response
import google.generativeai as genai
import json
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE 
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from dotenv import load_dotenv
import os
import io
import requests
from pptx.oxml import parse_xml
app = Flask(__name__)

# Load the environment variables from the .env file
load_dotenv()

# Get the API keys from the environment variables
api_key = os.getenv("API_KEY")
pexels_api_key = os.getenv("PEXELS_API_KEY")
genai.configure(api_key=api_key)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    presentation_title = request.form['title']
    number_of_slides = request.form['slides']

    # Define the prompt for generating the presentation slides.
    prompt = f"""
    generate a {number_of_slides} slide presentation for the topic {presentation_title}. Each slide should have a (header), (content). Return as JSON. Must include 3 bullet points.
    Using this JSON schema:
        {{
            "slides": [
                {{
                    "header": "string",
                    "content": "string"
                }}
            ]
        }}
    """

    # Choose a model that's appropriate for your use case.
    model = genai.GenerativeModel('gemini-1.5-flash',
                                  generation_config={"response_mime_type": "application/json"})

    # Send the request to the Gemini API.
    response = model.generate_content(prompt)

    # Parse the JSON response.
    response_data = json.loads(response.text)
    slide_data = response_data["slides"]

    # Create a PowerPoint presentation.
    prs = Presentation()

    for slide in slide_data:
        slide_layout = prs.slide_layouts[6]  # Use a blank slide layout without any text boxes
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Fetch relevant image using Pexels API
        image_url = None
        if slide['header']:
            query = slide['header']
            pexels_url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
            headers = {"Authorization": pexels_api_key}
            response = requests.get(pexels_url, headers=headers)
            if response.status_code == 200:
                image_data = response.json()
                if image_data['photos']:
                    image_url = image_data['photos'][0]['src']['medium']
                    
        # Add image to slide background
        if image_url:
            image_response = requests.get(image_url)
            image_stream = io.BytesIO(image_response.content)
            new_slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

            # Add a semi-transparent rectangle as text background
            left = Inches(1)
            top = Inches(1)
            width = prs.slide_width - Inches(2)
            height = prs.slide_height - Inches(2)

            shape = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )


            shape.fill.solid()
            fill_color = shape.fill.fore_color
            fill_color.rgb = RGBColor(255, 255, 255)

            # Modify XML to set transparency
            sp = shape._element
            solidFill = sp.find(".//a:solidFill", sp.nsmap)
            srgbClr = solidFill.find(".//a:srgbClr", sp.nsmap)

            # Add transparency (0.3 means 30% transparent)
            alpha = parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="40000"/>')  # 30% transparent
            srgbClr.append(alpha)

            shape.line.fill.background()  # Make the border transparent


        # Add the header and content text on top of the rectangle
        text_box = new_slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # Add the header
        if slide['header']:
            p = text_frame.add_paragraph()
            p.text = slide['header']
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
            p.space_after = Pt(14)

        # Add the content as bullet points
        if slide['content']:
            content_lines = slide['content'].split('\n')
            for line in content_lines:
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black font color
                p.level = 1  # Make it a bullet point

    # Save the PowerPoint presentation to a BytesIO object.
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)

    # Create a Flask response with the PowerPoint file.
    response = make_response(pptx_file.read())
    response.headers.set('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response.headers.set('Content-Disposition', 'attachment', filename='presentation.pptx')

    return response
